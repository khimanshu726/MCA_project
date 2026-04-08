from pathlib import Path
from tempfile import NamedTemporaryFile
from urllib.request import urlopen

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "Elite_Empressions_Project_Presentation_Final.pptx"
W = Inches(13.333)
H = Inches(7.5)

ORDER_DIAGRAM_URL = "https://s3-alpha.figma.com/thumbnails/80ac52d6-a5a4-498c-96ee-de9ae50c1e9f?X-Amz-Algorithm=AWS4-HMAC-SHA256&X-Amz-Credential=AKIAQ4GOSFWCYKED6IIG%2F20260408%2Fus-west-2%2Fs3%2Faws4_request&X-Amz-Date=20260408T194929Z&X-Amz-Expires=604800&X-Amz-SignedHeaders=host&X-Amz-Signature=0ace8bf5bedae1b5a40fc428957bc34fb441558480d16ec705f6afc5efdb156d"
ARCH_DIAGRAM_URL = "https://s3-alpha.figma.com/thumbnails/14fe51df-fc05-4ee6-9f35-0e661b656bbd?X-Amz-Algorithm=AWS4-HMAC-SHA256&X-Amz-Credential=AKIAQ4GOSFWCYKED6IIG%2F20260408%2Fus-west-2%2Fs3%2Faws4_request&X-Amz-Date=20260408T194931Z&X-Amz-Expires=604800&X-Amz-SignedHeaders=host&X-Amz-Signature=d373377cbb240e0922a59c50056ea65811ed8059f32c9d42f0e4d0be50555b69"

try:
    from cairosvg import svg2png
except Exception:
    svg2png = None

C = {
    "navy": RGBColor(14, 23, 44),
    "slate": RGBColor(46, 62, 92),
    "ink": RGBColor(26, 33, 53),
    "teal": RGBColor(35, 161, 166),
    "mint": RGBColor(91, 203, 171),
    "amber": RGBColor(232, 146, 55),
    "rose": RGBColor(204, 92, 92),
    "sky": RGBColor(102, 166, 255),
    "bg": RGBColor(244, 247, 251),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(216, 223, 232),
    "text": RGBColor(31, 41, 55),
    "muted": RGBColor(92, 108, 128),
    "white": RGBColor(255, 255, 255),
}


def box(slide, x, y, w, h, text="", size=20, color=None, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color or C["text"]
    p.alignment = align
    return shape


def bullets(shape, items, size=16, color=None, font="Aptos"):
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.font.size = Pt(size)
        p.font.name = font
        p.font.color.rgb = color or C["text"]


def card(slide, x, y, w, h, fill=None, line=None, shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill or C["card"]
    s.line.color.rgb = line or C["line"]
    return s


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def header(slide, title, sub=""):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C["teal"]
    bar.line.fill.background()
    box(slide, Inches(0.7), Inches(0.55), Inches(9), Inches(0.45), title, 28, C["ink"], True)
    if sub:
        box(slide, Inches(0.72), Inches(1.08), Inches(11.3), Inches(0.25), sub, 12, C["muted"])


def footer(slide, text):
    box(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.2), text, 10, C["muted"])


def two_col_slide(prs, title, sub, left_title, left_items, right_title, right_items, foot):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, C["bg"])
    header(slide, title, sub)
    card(slide, Inches(0.7), Inches(1.6), Inches(5.95), Inches(4.95))
    card(slide, Inches(6.95), Inches(1.6), Inches(5.7), Inches(4.95))
    box(slide, Inches(1.0), Inches(1.9), Inches(5), Inches(0.3), left_title, 18, C["ink"], True)
    bullets(box(slide, Inches(0.98), Inches(2.3), Inches(5.2), Inches(3.9)), left_items, 16)
    box(slide, Inches(7.22), Inches(1.9), Inches(5), Inches(0.3), right_title, 18, C["ink"], True)
    bullets(box(slide, Inches(7.18), Inches(2.3), Inches(5), Inches(3.9)), right_items, 16)
    footer(slide, foot)
    return slide


def code_panel(slide, x, y, w, h, color, title_text, content):
    card(slide, x, y, w, h, C["card"], C["card"])
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.38))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    box(slide, x + Inches(0.12), y + Inches(0.46), w - Inches(0.24), Inches(0.24), title_text, 16, C["ink"], True)
    box(slide, x + Inches(0.15), y + Inches(0.92), w - Inches(0.3), h - Inches(1.1), content, 11, C["text"], False, font="Courier New")


def add_remote_picture(slide, url, x, y, w, h):
    with urlopen(url) as response:
        content_type = response.headers.get("Content-Type", "")
        image_bytes = response.read()

    if "svg" in content_type:
        if svg2png is None:
            return False
        image_bytes = svg2png(bytestring=image_bytes)

    with NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
        temp_file.write(image_bytes)
        temp_path = temp_file.name

    slide.shapes.add_picture(temp_path, x, y, w, h)
    return True


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title = prs.slides.add_slide(prs.slide_layouts[6])
    bg(title, C["navy"])
    band = title.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.22))
    band.fill.solid()
    band.fill.fore_color.rgb = C["amber"]
    band.line.fill.background()
    pill = title.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(1.6), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = C["teal"]
    pill.line.fill.background()
    box(title, Inches(0.8), Inches(0.88), Inches(1.6), Inches(0.2), "PROJECT PPT", 12, C["white"], True, PP_ALIGN.CENTER)
    box(title, Inches(0.8), Inches(1.45), Inches(7.5), Inches(0.7), "Elite Empressions", 30, C["white"], True)
    box(title, Inches(0.8), Inches(2.2), Inches(8.8), Inches(0.7), "Printing Press E-Commerce & Order Management System", 22, C["white"], True)
    box(title, Inches(0.8), Inches(3.0), Inches(9.8), Inches(0.5), "Prepared according to CAPOL601R01 Project Work & Viva Voce Guidelines", 16, C["bg"])
    box(title, Inches(0.8), Inches(3.75), Inches(9.5), Inches(0.4), "React + Vite + Express + JWT + Passport + Order APIs", 14, C["mint"])
    card(title, Inches(8.9), Inches(1.05), Inches(3.55), Inches(4.9), C["card"], C["card"])
    box(title, Inches(9.18), Inches(1.35), Inches(3), Inches(0.3), "Presentation Highlights", 18, C["ink"], True)
    bullets(
        box(title, Inches(9.12), Inches(1.85), Inches(3), Inches(2.9)),
        [
            "Problem statement and abstract",
            "Project plan and SRS",
            "System analysis and design",
            "Coding, testing, implementation",
            "Future plans and conclusion",
        ],
        15,
    )
    footer(title, "Elite Empressions | Project Work Viva Voce Presentation")

    two_col_slide(
        prs,
        "Introduction",
        "Overview of the application and its purpose.",
        "Project Overview",
        [
            "Elite Empressions is a print-shop web application for browsing, customizing, and ordering print products.",
            "The system combines customer ordering and admin-side order management in one connected workflow.",
            "Products include visiting cards, flyers, packaging, T-shirts, mugs, invitations, stationery, and banners.",
        ],
        "Why This Project Matters",
        [
            "Traditional print ordering is often manual, fragmented, and slow.",
            "Customers need preview, customization, pricing clarity, and smoother checkout.",
            "Admins need one place to manage orders, files, statuses, and customer details.",
        ],
        "Assignment 2 | Introduction and project context",
    )

    two_col_slide(
        prs,
        "Problem Statement",
        "The current printing press workflow is inefficient for both customers and administrators.",
        "Existing Gaps",
        [
            "Customers cannot always preview or personalize products before placing orders.",
            "Design files, delivery details, and payment notes are often collected inconsistently.",
            "Admins spend extra time tracking order status and production information manually.",
        ],
        "Proposed Solution",
        [
            "A responsive storefront with guided product discovery and customization.",
            "A unified cart and checkout flow for addresses, payments, and artwork uploads.",
            "A secure admin dashboard to review, manage, and update the order lifecycle.",
        ],
        "Assignment 2 | Problem definition",
    )

    two_col_slide(
        prs,
        "Abstract",
        "Summary of the implemented application and its academic scope.",
        "Abstract",
        [
            "The project delivers a web-based order management system for a printing press using React, Vite, and Express.",
            "Customers can browse products, upload designs, manage cart items, save addresses, and place orders.",
            "Admins can authenticate, view incoming orders, download uploads, and update order status.",
            "The system demonstrates practical integration of frontend, backend, validation, persistence, and authentication modules.",
        ],
        "Scope",
        [
            "Focuses on the core e-commerce workflow for a print service business.",
            "Uses lightweight local storage and JSON-backed services for demonstration.",
            "Designed so the current prototype can be extended into a production-ready solution later.",
        ],
        "Assignment 1 | Project title and abstract",
    )

    plan = prs.slides.add_slide(prs.slide_layouts[6])
    bg(plan, C["bg"])
    header(plan, "Software Project Plan", "Development roadmap aligned to the review points in the guideline PDF.")
    phases = [
        ("Phase 1", "Requirement study, project title, abstract, and problem analysis", C["teal"]),
        ("Phase 2", "UI routing, catalog, customization module, cart, and checkout", C["mint"]),
        ("Phase 3", "Order APIs, admin panel, auth flows, uploads, and persistence", C["sky"]),
        ("Phase 4", "Testing, debugging, documentation, and viva presentation", C["amber"]),
    ]
    for i, (name, text, color) in enumerate(phases):
        x = Inches(0.9 + i * 3.05)
        step = plan.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Inches(2.7), Inches(0.72), Inches(0.72))
        step.fill.solid()
        step.fill.fore_color.rgb = color
        step.line.fill.background()
        box(plan, x, Inches(2.86), Inches(0.72), Inches(0.2), str(i + 1), 16, C["white"], True, PP_ALIGN.CENTER)
        card(plan, x - Inches(0.2), Inches(1.55), Inches(2.1), Inches(0.95))
        box(plan, x - Inches(0.08), Inches(1.83), Inches(1.85), Inches(0.2), name, 16, C["ink"], True, PP_ALIGN.CENTER)
        card(plan, x - Inches(0.45), Inches(3.8), Inches(2.6), Inches(1.55))
        box(plan, x - Inches(0.28), Inches(4.1), Inches(2.25), Inches(1.0), text, 13, C["text"], False, PP_ALIGN.CENTER)
    footer(plan, "Review 1: analysis and planning | Review 2: design, coding, testing, implementation")

    two_col_slide(
        prs,
        "Software Requirement Specification",
        "Core functional and non-functional requirements of the application.",
        "Functional Requirements",
        [
            "Display product catalog with categories, images, prices, and detail pages.",
            "Support customization preview, file upload, and order notes.",
            "Maintain shared cart state with add, update, remove, and persistence.",
            "Collect delivery details with address validation and city autocomplete.",
            "Create orders and allow admin-side tracking and status updates.",
            "Authenticate admins through email/password, OTP, and Google sign-in.",
        ],
        "Non-Functional Requirements",
        [
            "Responsive design for desktop and mobile.",
            "Secure password hashing and JWT-based route protection.",
            "Predictable cart and checkout behavior with real-time validation.",
            "Reusable components and modular project structure.",
            "Accessible labels, alt text, and professional visual design.",
            "Fast local development and easy future scalability.",
        ],
        "Assignment 2 | SRS",
    )

    analysis = prs.slides.add_slide(prs.slide_layouts[6])
    bg(analysis, C["bg"])
    header(analysis, "System Analysis", "Use case view for customer ordering flow and admin order management.")
    actors = [
        (Inches(0.9), "Customer", C["ink"]),
        (Inches(10.75), "Admin", C["slate"]),
    ]
    for x, label, color in actors:
        card(analysis, x, Inches(2.35), Inches(1.7), Inches(1.2), color, color)
        box(analysis, x + Inches(0.12), Inches(2.72), Inches(1.45), Inches(0.25), label, 20, C["white"], True, PP_ALIGN.CENTER)
    cases = [
        (Inches(3.05), Inches(1.7), "Browse Products", C["teal"]),
        (Inches(5.1), Inches(1.7), "Customize Product", C["teal"]),
        (Inches(7.15), Inches(1.7), "Add to Cart", C["teal"]),
        (Inches(3.05), Inches(3.2), "Place Order", C["mint"]),
        (Inches(5.1), Inches(3.2), "Upload Design File", C["mint"]),
        (Inches(7.15), Inches(3.2), "Track Status", C["mint"]),
        (Inches(4.05), Inches(4.7), "Review Orders", C["amber"]),
        (Inches(6.2), Inches(4.7), "Update Order Stage", C["amber"]),
    ]
    for x, y, label, color in cases:
        card(analysis, x, y, Inches(1.85), Inches(0.72), color, color)
        box(analysis, x + Inches(0.08), y + Inches(0.2), Inches(1.68), Inches(0.2), label, 13, C["white"], True, PP_ALIGN.CENTER)
    footer(analysis, "Actors: Customer and Admin | Core artifacts: products, cart, orders, and uploaded files")

    order_flow = prs.slides.add_slide(prs.slide_layouts[6])
    bg(order_flow, C["bg"])
    header(order_flow, "Detailed Order Flow Diagram", "Expanded customer-to-admin lifecycle for the Elite Empressions application.")
    order_flow_image_added = add_remote_picture(order_flow, ORDER_DIAGRAM_URL, Inches(0.75), Inches(1.65), Inches(7.5), Inches(4.8))
    if not order_flow_image_added:
        flow_steps = [
            (Inches(0.95), Inches(2.2), "Browse\nProducts", C["teal"]),
            (Inches(2.45), Inches(2.2), "Product\nDetails", C["sky"]),
            (Inches(3.95), Inches(2.2), "Customize\nand Upload", C["mint"]),
            (Inches(5.45), Inches(2.2), "Cart and\nCheckout", C["amber"]),
            (Inches(1.7), Inches(4.05), "Create\nOrder API", C["rose"]),
            (Inches(3.7), Inches(4.05), "Admin\nDashboard", C["slate"]),
            (Inches(5.7), Inches(4.05), "Status\nUpdate", C["ink"]),
        ]
        for x, y, label, color in flow_steps:
            card(order_flow, x, y, Inches(1.2), Inches(0.9), color, color)
            box(order_flow, x + Inches(0.05), y + Inches(0.18), Inches(1.1), Inches(0.4), label, 13, C["white"], True, PP_ALIGN.CENTER)
        arrows = [
            (Inches(2.15), Inches(2.55)), (Inches(3.65), Inches(2.55)), (Inches(5.15), Inches(2.55)),
            (Inches(2.25), Inches(4.4)), (Inches(4.25), Inches(4.4)),
        ]
        for x, y in arrows:
            box(order_flow, x, y, Inches(0.25), Inches(0.2), "→", 22, C["muted"], True, PP_ALIGN.CENTER)
        box(order_flow, Inches(6.1), Inches(2.95), Inches(0.35), Inches(0.3), "↓", 22, C["muted"], True, PP_ALIGN.CENTER)
        box(order_flow, Inches(3.05), Inches(1.65), Inches(1.8), Inches(0.25), "Customer Journey", 16, C["ink"], True, PP_ALIGN.CENTER)
        box(order_flow, Inches(3.0), Inches(5.1), Inches(2.0), Inches(0.25), "Operations Journey", 16, C["ink"], True, PP_ALIGN.CENTER)
    card(order_flow, Inches(8.55), Inches(1.75), Inches(4.1), Inches(4.65))
    box(order_flow, Inches(8.85), Inches(2.0), Inches(3.4), Inches(0.3), "Diagram Highlights", 18, C["ink"], True)
    bullets(
        box(order_flow, Inches(8.75), Inches(2.45), Inches(3.35), Inches(2.05)),
        [
            "The diagram traces the journey from browsing products to final operational order closure.",
            "A clear handoff is shown between customer activity and admin-side processing.",
            "The flow matches the implemented checkout, order API, and admin dashboard lifecycle.",
        ],
        14,
    )
    bullets(
        box(order_flow, Inches(8.75), Inches(4.6), Inches(3.35), Inches(1.45)),
        [
            "This makes the viva explanation more concrete than a generic use-case sketch.",
            "The visual structure follows the real routing and processing behavior of the app.",
        ],
        14,
    )
    footer(order_flow, "Expanded diagram view for viva explanation of complete order processing")

    frontend = prs.slides.add_slide(prs.slide_layouts[6])
    bg(frontend, C["bg"])
    header(frontend, "Design - Front End", "Representative interface structure based on the implemented React application.")
    panels = [
        (Inches(0.9), "Home", C["teal"], ["Hero banner", "Category grid", "Quick navigation"]),
        (Inches(3.95), "Products", C["sky"], ["Product cards", "Detail routing", "Add to cart feedback"]),
        (Inches(7.0), "Cart + Checkout", C["mint"], ["Address form", "Payment options", "Order summary"]),
        (Inches(10.05), "Admin", C["amber"], ["Order table", "Filters", "Status controls"]),
    ]
    for x, title_text, color, items in panels:
        card(frontend, x, Inches(1.8), Inches(2.35), Inches(4.65))
        top = frontend.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.8), Inches(2.35), Inches(0.35))
        top.fill.solid()
        top.fill.fore_color.rgb = C["ink"]
        top.line.fill.background()
        card(frontend, x + Inches(0.18), Inches(2.28), Inches(1.95), Inches(0.55), color, color)
        box(frontend, x + Inches(0.24), Inches(2.45), Inches(1.8), Inches(0.2), title_text, 15, C["white"], True, PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            card(frontend, x + Inches(0.18), Inches(3.02 + j * 0.66), Inches(1.95), Inches(0.52), C["bg"], C["line"])
            box(frontend, x + Inches(0.26), Inches(3.16 + j * 0.66), Inches(1.78), Inches(0.2), item, 11, C["text"], False, PP_ALIGN.CENTER)
    footer(frontend, "React functional components, hooks, React Router, reusable UI cards, and shared context state")

    backend = prs.slides.add_slide(prs.slide_layouts[6])
    bg(backend, C["bg"])
    header(backend, "Design - Back End", "Server architecture and data flow of the Express-based order management system.")
    nodes = [
        (Inches(0.85), Inches(2.3), Inches(2.0), Inches(1.0), "React Client", C["teal"]),
        (Inches(3.45), Inches(1.35), Inches(2.3), Inches(1.0), "Auth Routes", C["sky"]),
        (Inches(3.45), Inches(3.25), Inches(2.3), Inches(1.0), "Order Routes", C["mint"]),
        (Inches(6.3), Inches(1.35), Inches(2.45), Inches(1.0), "Controllers", C["amber"]),
        (Inches(6.3), Inches(3.25), Inches(2.45), Inches(1.0), "Services", C["rose"]),
        (Inches(9.35), Inches(1.35), Inches(2.55), Inches(1.0), "users.json / orders.json", C["ink"]),
        (Inches(9.35), Inches(3.25), Inches(2.55), Inches(1.0), "Uploads Folder", C["slate"]),
    ]
    for x, y, w, h, label, color in nodes:
        card(backend, x, y, w, h, color, color)
        box(backend, x + Inches(0.08), y + Inches(0.22), w - Inches(0.16), Inches(0.45), label, 14, C["white"], True, PP_ALIGN.CENTER)
    box(backend, Inches(0.95), Inches(5.35), Inches(11.3), Inches(0.6), "The backend exposes REST APIs, validates payloads, stores development data locally, handles file uploads with Multer, and protects admin endpoints using JWT middleware.", 14, C["text"], False, PP_ALIGN.CENTER)
    footer(backend, "Node.js + Express 5 + Passport Google OAuth + bcryptjs + Multer + Nodemailer")

    api_arch = prs.slides.add_slide(prs.slide_layouts[6])
    bg(api_arch, C["bg"])
    header(api_arch, "Detailed API and Architecture Diagram", "Layered view of the frontend, API, service, and storage interactions.")
    api_arch_image_added = add_remote_picture(api_arch, ARCH_DIAGRAM_URL, Inches(0.75), Inches(1.65), Inches(7.5), Inches(4.8))
    if not api_arch_image_added:
        layers = [
            (Inches(0.95), "Frontend", ["Home", "Products", "Customize", "Cart", "Admin"], C["teal"]),
            (Inches(2.85), "API", ["/api/auth", "/api/orders", "/api/admin"], C["sky"]),
            (Inches(4.75), "Controllers", ["Auth", "Orders", "Validation"], C["amber"]),
            (Inches(6.65), "Services", ["User Store", "Order Store", "OTP Store", "Uploads"], C["rose"]),
        ]
        for x, title_text, items, color in layers:
            card(api_arch, x, Inches(2.1), Inches(1.45), Inches(3.1), color, color)
            box(api_arch, x + Inches(0.08), Inches(2.28), Inches(1.28), Inches(0.22), title_text, 15, C["white"], True, PP_ALIGN.CENTER)
            for idx, item in enumerate(items):
                box(api_arch, x + Inches(0.12), Inches(2.68 + idx * 0.45), Inches(1.2), Inches(0.18), item, 10, C["white"], False, PP_ALIGN.CENTER)
        card(api_arch, Inches(0.95), Inches(5.45), Inches(2.5), Inches(0.7), C["slate"], C["slate"])
        card(api_arch, Inches(3.65), Inches(5.45), Inches(2.5), Inches(0.7), C["slate"], C["slate"])
        card(api_arch, Inches(6.35), Inches(5.45), Inches(1.75), Inches(0.7), C["mint"], C["mint"])
        box(api_arch, Inches(1.1), Inches(5.67), Inches(2.2), Inches(0.2), "users.json", 13, C["white"], True, PP_ALIGN.CENTER)
        box(api_arch, Inches(3.8), Inches(5.67), Inches(2.2), Inches(0.2), "orders.json", 13, C["white"], True, PP_ALIGN.CENTER)
        box(api_arch, Inches(6.48), Inches(5.67), Inches(1.5), Inches(0.2), "uploads", 13, C["white"], True, PP_ALIGN.CENTER)
        for x in [Inches(2.45), Inches(4.35), Inches(6.25)]:
            box(api_arch, x, Inches(3.45), Inches(0.25), Inches(0.2), "→", 22, C["muted"], True, PP_ALIGN.CENTER)
        box(api_arch, Inches(6.95), Inches(4.9), Inches(0.3), Inches(0.2), "↓", 22, C["muted"], True, PP_ALIGN.CENTER)
    card(api_arch, Inches(8.55), Inches(1.75), Inches(4.1), Inches(4.65))
    box(api_arch, Inches(8.85), Inches(2.0), Inches(3.4), Inches(0.3), "Architecture Highlights", 18, C["ink"], True)
    bullets(
        box(api_arch, Inches(8.75), Inches(2.45), Inches(3.35), Inches(2.05)),
        [
            "The layered architecture separates frontend screens, API routes, controller logic, and persistence services.",
            "Auth and order concerns are split cleanly, improving maintainability and debugging clarity.",
            "The diagram reflects the actual modules used in the Express backend and React frontend.",
        ],
        14,
    )
    bullets(
        box(api_arch, Inches(8.75), Inches(4.6), Inches(3.35), Inches(1.45)),
        [
            "This structure is easy to explain in viva and maps directly to implemented files and routes.",
            "It also shows readiness for later migration to a production database.",
        ],
        14,
    )
    footer(api_arch, "Layered architecture slide for backend explanation in viva")

    two_col_slide(
        prs,
        "Interface Design",
        "Major customer and admin interfaces implemented in the project.",
        "Customer Interfaces",
        [
            "Home page with hero banner and category entry cards",
            "Product list and detail pages with large images and gallery view",
            "Customization page with upload preview and overlay support",
            "Cart + checkout page with address book, payment, and order creation",
        ],
        "Admin Interfaces",
        [
            "Login and registration with email, mobile OTP, and Google options",
            "Protected order dashboard with filters and metrics",
            "Detailed order panel with line items, payment state, and file download",
            "Status controls for New, Processing, Completed, and Cancelled orders",
        ],
        "Assignment 3 | Interface design",
    )

    data_model = prs.slides.add_slide(prs.slide_layouts[6])
    bg(data_model, C["bg"])
    header(data_model, "Back End Design - Data Model", "Logical structure of user, order, and OTP records.")
    for x, title_text, items in [
        (Inches(0.85), "User Record", ["id", "email", "mobile", "password (hashed)", "provider", "role", "createdAt"]),
        (Inches(4.95), "Order Record", ["orderId", "customerName, phone, email", "address fields", "productName, quantity, price", "uploadedFileURL", "paymentStatus, orderStatus", "lineItems, createdAt"]),
        (Inches(9.45), "OTP Store", ["mobile", "otp", "expiresAt", "request throttle", "temporary memory storage"]),
    ]:
        width = Inches(3.7 if x < Inches(1) else 4.15 if x < Inches(6) else 2.95)
        card(data_model, x, Inches(1.75), width, Inches(4.9))
        box(data_model, x + Inches(0.25), Inches(2.0), width - Inches(0.5), Inches(0.3), title_text, 18, C["ink"], True)
        bullets(box(data_model, x + Inches(0.2), Inches(2.4), width - Inches(0.4), Inches(3.4)), items, 14)
    footer(data_model, "Current implementation uses local persistence and is ready for future database migration")

    cart_code = prs.slides.add_slide(prs.slide_layouts[6])
    bg(cart_code, C["navy"])
    box(cart_code, Inches(0.7), Inches(0.55), Inches(5.6), Inches(0.5), "Sample Code - Cart State Management", 28, C["white"], True)
    box(cart_code, Inches(0.72), Inches(1.05), Inches(9.4), Inches(0.3), "Focused code snippet slide for the centralized cart logic used across listing, detail, and cart pages.", 13, C["bg"])
    code_panel(
        cart_code,
        Inches(0.7),
        Inches(1.65),
        Inches(6.2),
        Inches(4.9),
        C["teal"],
        "Reducer-based cart flow",
        "function addToCart(product) {\n  if (item exists) {\n    quantity = quantity + 1;\n  } else {\n    quantity = 1;\n  }\n}\n\nfunction updateQuantity(id, action) {\n  if (action === 'increment') quantity += 1;\n  if (action === 'decrement' && quantity === 1) remove item;\n  else quantity -= 1;\n}",
    )
    bullets(
        box(cart_code, Inches(7.2), Inches(1.82), Inches(5.0), Inches(2.3)),
        [
            "Centralized Context API plus reducer removed duplicate cart logic across pages.",
            "Quantity changes are deterministic and always happen by plus or minus one.",
            "localStorage persistence keeps the cart intact after refresh.",
            "This solved the earlier inconsistent add-to-cart behavior.",
        ],
        15,
        C["white"],
    )
    footer(cart_code, "Assignment 3 | Code snippet focused on reliable cart behavior")

    auth_code = prs.slides.add_slide(prs.slide_layouts[6])
    bg(auth_code, C["navy"])
    box(auth_code, Inches(0.7), Inches(0.55), Inches(5.6), Inches(0.5), "Sample Code - Authentication and OTP", 28, C["white"], True)
    box(auth_code, Inches(0.72), Inches(1.05), Inches(9.4), Inches(0.3), "Focused code snippet slide for email login, mobile OTP, and backend verification flow.", 13, C["bg"])
    code_panel(
        auth_code,
        Inches(0.7),
        Inches(1.65),
        Inches(5.8),
        Inches(4.9),
        C["amber"],
        "Backend OTP handlers",
        "POST /auth/send-otp\nvalidate mobile\notp = random 6 digits\nsaveOtpForMobile(mobile, otp)\nsendOtpSms(mobile, otp)\n\nPOST /auth/verify-otp\nverifyOtpForMobile(mobile, otp)\nfind or create user\nissue JWT token",
    )
    code_panel(
        auth_code,
        Inches(6.75),
        Inches(1.65),
        Inches(5.55),
        Inches(4.9),
        C["mint"],
        "Frontend detection flow",
        "if (identifier includes '@') {\n  loginType = 'email';\n  show password field;\n} else if (identifier is numeric) {\n  loginType = 'mobile';\n  show OTP flow;\n}\n\nGoogle button -> /api/auth/google",
    )
    footer(auth_code, "Assignment 3 | Code snippets for authentication logic")

    order_code = prs.slides.add_slide(prs.slide_layouts[6])
    bg(order_code, C["navy"])
    box(order_code, Inches(0.7), Inches(0.55), Inches(5.6), Inches(0.5), "Sample Code - Order Processing", 28, C["white"], True)
    box(order_code, Inches(0.72), Inches(1.05), Inches(9.4), Inches(0.3), "Focused code snippet slide for the checkout-to-order backend flow and status lifecycle.", 13, C["bg"])
    code_panel(
        order_code,
        Inches(0.7),
        Inches(1.65),
        Inches(6.2),
        Inches(4.9),
        C["rose"],
        "POST /api/orders",
        "parse lineItems\nvalidate customer + address + payment\ncalculate quantity and total price\nset paymentStatus = Pending or Paid\ncreate uploadedFileURL\ncreate order record\nsend admin notification\nreturn created order",
    )
    bullets(
        box(order_code, Inches(7.2), Inches(1.82), Inches(5.0), Inches(2.3)),
        [
            "Order lifecycle used in admin panel: New -> Processing -> Completed or Cancelled.",
            "Uploaded PDF, PNG, and JPG files are linked with the order for admin access.",
            "Payment status is derived from selected payment method during checkout.",
            "This code connects customer checkout directly to business operations.",
        ],
        15,
        C["white"],
    )
    footer(order_code, "Assignment 3 | Code snippet focused on order creation and status management")

    two_col_slide(
        prs,
        "Testing",
        "Validation and debugging coverage considered during development.",
        "Testing Types",
        [
            "Unit testing of address validation, OTP payload checks, and cart quantity logic",
            "Integration testing of checkout to order API, admin fetch, and order status updates",
            "Form validation for name, phone, email, pincode, and file uploads",
        ],
        "Debugging Highlights",
        [
            "Fixed inconsistent cart state and quantity increments across pages",
            "Resolved address persistence issues after delete and refresh",
            "Corrected Google OAuth and OTP route wiring and stale-server conflicts",
        ],
        "Assignment 3 | Testing and validation",
    )

    two_col_slide(
        prs,
        "Implementation, Problems Faced, and Lessons Learnt",
        "Key outcomes and learnings collected during development.",
        "Implementation Outcome",
        [
            "Built a working storefront with shared cart, checkout, and order placement flow",
            "Integrated backend APIs for order creation, uploads, authentication, and admin management",
            "Delivered a modular project structure with reusable components and maintainable logic",
        ],
        "Problems and Lessons",
        [
            "State must be centralized to prevent duplicate cart logic and inconsistent behavior",
            "Clear server logs and route verification reduce time spent debugging auth issues",
            "Validation and persistence logic are as important as the visible UI in e-commerce flows",
        ],
        "Assignment 3 | Implementation review",
    )

    two_col_slide(
        prs,
        "Future Plans",
        "Enhancements planned beyond the current implemented scope.",
        "Enhancement Areas",
        [
            "Migrate local JSON stores to MongoDB or PostgreSQL",
            "Integrate Razorpay or Stripe for real payments",
            "Add advanced design editor tools and saved templates",
            "Introduce analytics, invoice generation, and dispatch integration",
        ],
        "Expected Benefits",
        [
            "Better scalability and production readiness",
            "Improved customer trust with secure online payments",
            "Stronger personalization and conversion rates",
            "More efficient day-to-day operations for the print business",
        ],
        "Assignment 3 | Future plans",
    )

    end = prs.slides.add_slide(prs.slide_layouts[6])
    bg(end, C["navy"])
    box(end, Inches(0.85), Inches(1.05), Inches(6.2), Inches(0.6), "Conclusion", 30, C["white"], True)
    box(end, Inches(0.85), Inches(1.85), Inches(7.0), Inches(2.2), "Elite Empressions demonstrates a complete digital workflow for a modern printing press: product discovery, customization, cart and checkout, order creation, file handling, and admin-side order tracking.", 20, C["bg"])
    box(end, Inches(0.85), Inches(4.25), Inches(7.0), Inches(1.0), "The project satisfies the guideline structure for analysis, design, coding, testing, implementation, future scope, and final presentation.", 18, C["mint"])
    card(end, Inches(8.55), Inches(1.25), Inches(3.9), Inches(4.85), C["card"], C["card"])
    box(end, Inches(8.82), Inches(1.55), Inches(3.2), Inches(0.3), "References", 18, C["ink"], True)
    bullets(
        box(end, Inches(8.78), Inches(2.0), Inches(3.0), Inches(2.0)),
        [
            "CAPOL601R01 Project Work & Viva Voce Guidelines PDF",
            "Elite Empressions source code and application modules",
            "React, Vite, Express, Passport, and JWT documentation",
        ],
        15,
    )
    box(end, Inches(8.9), Inches(5.0), Inches(2.8), Inches(0.5), "Thank you", 24, C["teal"], True, PP_ALIGN.CENTER)
    footer(end, "Prepared for project work viva voce review")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
